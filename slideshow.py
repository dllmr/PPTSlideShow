# /// script
# requires-python = ">=3.11"
# dependencies = [
#     "python-pptx>=1.0",
#     "Pillow>=10.0",
#     "lxml>=4.9",
# ]
# ///
"""Build a looping slideshow.pptx from images found recursively in the cwd.

Images are linked (not embedded) with paths relative to the PPTX, so the
PPTX must live at the root of the image tree when opened.
"""

import argparse
import sys
import tomllib
from io import BytesIO
from pathlib import Path
from urllib.parse import quote

from lxml import etree
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn
from pptx.util import Emu

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tif", ".tiff", ".webp"}

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
P14_NS = "http://schemas.microsoft.com/office/powerpoint/2010/main"

FADE_MS = 500

OUT_NAME = "slideshow.pptx"
CONFIG_NAME = "slideshow.toml"

DEFAULTS = {"duration": 10.0, "loop": True, "embed": False, "fade": False}


def load_config(path: Path) -> dict:
    cfg = dict(DEFAULTS)
    if not path.is_file():
        return cfg
    try:
        with path.open("rb") as f:
            loaded = tomllib.load(f)
    except (OSError, tomllib.TOMLDecodeError) as e:
        print(f"Warning: could not read {path.name}: {e}", file=sys.stderr)
        return cfg
    for key in DEFAULTS:
        if key in loaded:
            cfg[key] = loaded[key]
    if not isinstance(cfg["duration"], (int, float)) or cfg["duration"] <= 0:
        cfg["duration"] = DEFAULTS["duration"]
    for key in ("loop", "embed", "fade"):
        if not isinstance(cfg[key], bool):
            cfg[key] = DEFAULTS[key]
    return cfg


def save_config(path: Path, duration: float, loop: bool,
                embed: bool, fade: bool) -> None:
    body = (
        f"duration = {duration}\n"
        f"loop = {'true' if loop else 'false'}\n"
        f"embed = {'true' if embed else 'false'}\n"
        f"fade = {'true' if fade else 'false'}\n"
    )
    path.write_text(body)


def find_images(root: Path, exclude: Path) -> list[Path]:
    results = []
    for p in root.rglob("*"):
        if not p.is_file():
            continue
        if p.suffix.lower() not in IMAGE_EXTS:
            continue
        if p.resolve() == exclude.resolve():
            continue
        results.append(p)
    results.sort(key=lambda q: str(q.relative_to(root)).lower())
    return results


def set_black_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)


def add_linked_picture(slide, rel_url: str, left: int, top: int,
                       width: int, height: int, shape_id: int, name: str) -> None:
    rId = slide.part.relate_to(rel_url, RT.IMAGE, is_external=True)
    pic_xml = (
        f'<p:pic xmlns:a="{A_NS}" xmlns:p="{P_NS}" xmlns:r="{R_NS}">'
        f"<p:nvPicPr>"
        f'<p:cNvPr id="{shape_id}" name="{name}"/>'
        f'<p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>'
        f"<p:nvPr/>"
        f"</p:nvPicPr>"
        f"<p:blipFill>"
        f'<a:blip r:link="{rId}"/>'
        f"<a:stretch><a:fillRect/></a:stretch>"
        f"</p:blipFill>"
        f"<p:spPr>"
        f"<a:xfrm>"
        f'<a:off x="{left}" y="{top}"/>'
        f'<a:ext cx="{width}" cy="{height}"/>'
        f"</a:xfrm>"
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        f"</p:spPr>"
        f"</p:pic>"
    )
    slide.shapes._spTree.append(etree.fromstring(pic_xml))


def set_slide_auto_advance(slide, seconds: float, fade: bool) -> None:
    ms = max(1, int(round(seconds * 1000)))
    # The child element is load-bearing for LibreOffice: its importer only
    # applies advTm to the slide if the transition has a child element (or a
    # p14:dur attribute). Without it, LO parses advTm then silently drops it.
    effect = "fade" if fade else "cut"
    # p14:dur pins fade duration to exactly FADE_MS in modern readers;
    # spd="fast" is the legacy fallback (~0.5s) for readers that ignore p14.
    xml = (
        f'<p:transition xmlns:p="{P_NS}" xmlns:p14="{P14_NS}" '
        f'spd="fast" p14:dur="{FADE_MS}" advClick="0" advTm="{ms}">'
        f'<p:{effect}/>'
        f'</p:transition>'
    )
    transition = etree.fromstring(xml)
    sld = slide._element
    insert_idx = len(sld)
    for i, child in enumerate(sld):
        if child.tag == qn("p:timing"):
            insert_idx = i
            break
    sld.insert(insert_idx, transition)


def set_loop(prs) -> None:
    # LibreOffice reads loop from ppt/presProps.xml (<p:presentationPr>),
    # NOT from ppt/presentation.xml. See oox/source/ppt/presPropsfragmenthandler.cxx.
    # We set it on both for belt-and-braces coverage.
    ns = {"p": P_NS}

    pres = prs.part._element
    showPr = pres.find("p:showPr", ns)
    if showPr is None:
        before = {qn(f"p:{t}") for t in (
            "sldMasterIdLst", "notesMasterIdLst", "handoutMasterIdLst",
            "sldIdLst", "sldSz", "notesSz", "smartTags",
            "embeddedFontLst", "custShowLst",
        )}
        insert_idx = 0
        for i, child in enumerate(pres):
            if child.tag in before:
                insert_idx = i + 1
        showPr = etree.SubElement(pres, qn("p:showPr"))
        pres.remove(showPr)
        pres.insert(insert_idx, showPr)
    showPr.set("loop", "1")
    showPr.set("useTimings", "1")
    if showPr.find("p:sldAll", ns) is None:
        etree.SubElement(showPr, qn("p:sldAll"))

    props_part = next(
        (p for p in prs.part.package.iter_parts()
         if str(p.partname) == "/ppt/presProps.xml"),
        None,
    )
    if props_part is None:
        return
    props_root = etree.fromstring(props_part.blob)
    showPr2 = props_root.find("p:showPr", ns)
    if showPr2 is None:
        showPr2 = etree.Element(qn("p:showPr"))
        extLst = props_root.find("p:extLst", ns)
        if extLst is not None:
            extLst.addprevious(showPr2)
        else:
            props_root.append(showPr2)
    showPr2.set("loop", "1")
    showPr2.set("useTimings", "1")
    if showPr2.find("p:sldAll", ns) is None:
        sldAll = etree.Element(qn("p:sldAll"))
        showPr2.insert(0, sldAll)
    props_part._blob = etree.tostring(
        props_root, xml_declaration=True, encoding="UTF-8", standalone=True
    )


def build_thumbnail(first_img: Path, width: int = 320, height: int = 180) -> bytes:
    canvas = Image.new("RGB", (width, height), "black")
    with Image.open(first_img) as im:
        im = im.convert("RGB")
        scale = min(width / im.width, height / im.height)
        new_size = (max(1, int(im.width * scale)), max(1, int(im.height * scale)))
        im = im.resize(new_size, Image.LANCZOS)
        canvas.paste(im, ((width - new_size[0]) // 2, (height - new_size[1]) // 2))
    buf = BytesIO()
    canvas.save(buf, format="JPEG", quality=85)
    return buf.getvalue()


def set_thumbnail(prs, first_img: Path) -> None:
    for part in prs.part.package.iter_parts():
        if str(part.partname) == "/docProps/thumbnail.jpeg":
            try:
                part._blob = build_thumbnail(first_img)
            except Exception as e:
                print(f"Warning: could not build thumbnail: {e}", file=sys.stderr)
            return


def compute_fit(img_w: int, img_h: int,
                slide_w: int, slide_h: int) -> tuple[int, int, int, int]:
    scale = min(slide_w / img_w, slide_h / img_h)
    w = int(img_w * scale)
    h = int(img_h * scale)
    return (slide_w - w) // 2, (slide_h - h) // 2, w, h


def _yesno(prompt: str, default: bool) -> bool:
    hint = "[Y/n]" if default else "[y/N]"
    raw = input(f"{prompt} {hint}: ").strip().lower()
    if not raw:
        return default
    return raw in ("y", "yes")


def prompt_interactive(defaults: dict) -> tuple[float, bool, bool, bool]:
    d_dur = defaults["duration"]
    dur_str = f"{d_dur:g}"
    while True:
        raw = input(f"Slide duration in seconds [{dur_str}]: ").strip() or dur_str
        try:
            secs = float(raw)
            if secs > 0:
                break
        except ValueError:
            pass
        print("Please enter a positive number.")
    loop = _yesno("Loop the presentation?", defaults["loop"])
    embed = _yesno("Embed images instead of linking?", defaults["embed"])
    fade = _yesno("Fade transition between slides?", defaults["fade"])
    return secs, loop, embed, fade


def main() -> int:
    ap = argparse.ArgumentParser(
        description="Build slideshow.pptx from images under the current folder (recursive).",
    )
    ap.add_argument("-I", "--interactive", action="store_true",
                    help="Prompt for slide duration, loop, and embed/link.")
    args = ap.parse_args()

    root = Path.cwd()
    out_path = root / OUT_NAME
    config_path = root / CONFIG_NAME
    cfg = load_config(config_path)
    duration: float = float(cfg["duration"])
    loop: bool = cfg["loop"]
    embed: bool = cfg["embed"]
    fade: bool = cfg["fade"]
    if args.interactive:
        duration, loop, embed, fade = prompt_interactive(cfg)
        save_config(config_path, duration, loop, embed, fade)
    images = find_images(root, out_path)
    if not images:
        print("No images found.", file=sys.stderr)
        return 1

    prs = Presentation()
    prs.slide_width = Emu(12192000)   # 13.333"
    prs.slide_height = Emu(6858000)   # 7.5"
    slide_w, slide_h = prs.slide_width, prs.slide_height
    blank_layout = prs.slide_layouts[6]

    for idx, img_path in enumerate(images, start=1):
        slide = prs.slides.add_slide(blank_layout)
        set_black_background(slide)

        try:
            with Image.open(img_path) as im:
                iw, ih = im.size
        except Exception as e:
            print(f"Skipping {img_path}: {e}", file=sys.stderr)
            continue

        left, top, w, h = compute_fit(iw, ih, slide_w, slide_h)
        if embed:
            slide.shapes.add_picture(str(img_path), Emu(left), Emu(top),
                                     width=Emu(w), height=Emu(h))
        else:
            rel = img_path.relative_to(root).as_posix()
            rel_url = quote(rel)
            add_linked_picture(slide, rel_url, left, top, w, h,
                               shape_id=idx + 1, name=f"Picture {idx}")
        set_slide_auto_advance(slide, duration, fade)

    if loop:
        set_loop(prs)
    set_thumbnail(prs, images[0])

    prs.save(out_path)
    mode = "embedded" if embed else "linked"
    trans = "fade" if fade else "cut"
    print(f"Wrote {out_path.name} — {len(prs.slides)} slides, "
          f"{duration}s each, loop={loop}, images {mode}, transition {trans}.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
